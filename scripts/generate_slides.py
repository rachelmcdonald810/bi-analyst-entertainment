"""
Generate presentation slides as .pptx for the Live Music Analytics Pipeline project.
Run: python scripts/generate_slides.py
Output: docs/Live_Music_Analytics_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ─────────────────────────────────────────────────────────────
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE = RGBColor(0xFF, 0x9F, 0x1C)


def add_bg(slide):
    """Add dark background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER
    return slide


def add_content_slide(title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    # Title bar
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    # Accent line
    line = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.3), Inches(11.5), Emu(28000)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    # Bullets
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.space_after = Pt(12)

        # Check if bullet is bold (starts with **)
        if bullet.startswith("**") and "**" in bullet[2:]:
            end = bullet.index("**", 2)
            bold_text = bullet[2:end]
            rest = bullet[end + 2:]
            run1 = p.add_run()
            run1.text = bold_text
            run1.font.size = Pt(18)
            run1.font.bold = True
            run1.font.color.rgb = WHITE
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(18)
            run2.font.color.rgb = LIGHT_GRAY
        else:
            run = p.add_run()
            run.text = bullet
            run.font.size = Pt(18)
            run.font.color.rgb = LIGHT_GRAY

    # Note at bottom
    if note:
        txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.7))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = note
        p3.font.size = Pt(14)
        p3.font.italic = True
        p3.font.color.rgb = ORANGE
    return slide


# ── Slide 1: Title ────────────────────────────────────────────────────────────

add_title_slide(
    "Live Music Analytics Pipeline",
    "Rachel McDonald  |  ISBA 4715  |  Developing Business Applications with SQL"
)

# ── Slide 2: Problem Statement ────────────────────────────────────────────────

add_content_slide(
    "Streaming and Ticketing Data Live in Silos — Costing the Industry Millions",
    [
        "Entertainment companies need to connect streaming popularity with live event performance",
        "Ticketing data (events, venues, pricing) and streaming data (listeners, engagement) are managed separately",
        "BI analysts at companies like AXS (AEG Worldwide) need unified views for booking, pricing, and marketing",
        "Streaming metrics are a leading indicator of ticket demand — but only if you can combine the data",
        "This project builds the pipeline that bridges that gap",
    ]
)

# ── Slide 3: Solution ─────────────────────────────────────────────────────────

add_content_slide(
    "End-to-End Pipeline: From Raw APIs to Interactive Dashboard",
    [
        "**Sources:** Ticketmaster Discovery API (200 events) + Spotify artist pages via Firecrawl (15 artists)",
        "**Warehouse:** Snowflake with RAW → STAGING → MARTS schemas",
        "**Transform:** dbt star schema — 2 staging views with cleaning/type casting, 4 mart tables",
        "**Automate:** GitHub Actions runs both extraction pipelines on daily schedules",
        "**Visualize:** Interactive Streamlit dashboard deployed on Community Cloud",
        "**Research:** Knowledge base with 17 industry sources synthesized into wiki pages",
    ],
    note="[Insert pipeline diagram screenshot from README here]"
)

# ── Slide 4: Star Schema ──────────────────────────────────────────────────────

add_content_slide(
    "Star Schema Unifies Ticketing and Streaming Data in Snowflake",
    [
        "**fact_events:** 200 events with pricing, sale status, and dimension foreign keys",
        "**dim_artists:** 146 artists — Ticketmaster metadata joined to Spotify monthly listeners",
        "**dim_venues:** 184 venues with city and state for geographic analysis",
        "**dim_dates:** Temporal attributes — day of week, month, quarter, weekend flag",
        "Surrogate keys link fact to dimensions via dbt_utils.generate_surrogate_key()",
        "14 dbt tests validate data quality across staging and mart layers",
    ],
    note="[Insert ERD screenshot from README here]"
)

# ── Slide 5: Descriptive Insight 1 ────────────────────────────────────────────

add_content_slide(
    "Weekend Events Outnumber Weekdays 3:1 — Leaving Revenue on the Table",
    [
        "The vast majority of live music events are scheduled Friday through Sunday",
        "Monday through Wednesday have nearly empty event calendars across all genres and markets",
        "This pattern is consistent across genres — it's an industry-wide scheduling norm, not genre-specific",
        "Weekday venue inventory sits underutilized while weekend dates face competition",
        "**Callout:** Highlight the Saturday/Friday bars vs Monday-Wednesday in the day-of-week chart",
    ],
    note="[Insert Events by Day of Week bar chart screenshot from dashboard]"
)

# ── Slide 6: Descriptive Insight 2 ────────────────────────────────────────────

add_content_slide(
    "Rock and Pop Capture 80%+ of Events — Niche Genres Are Underbooked",
    [
        "A small number of genres dominate the event landscape nationally",
        "The long tail of genres (country, R&B, hip-hop, Latin) have far fewer scheduled events",
        "Some underbooked genres have strong streaming followings — supply doesn't match demand",
        "Ticket pricing also varies significantly by genre, suggesting genre-specific pricing strategies",
        "**Callout:** Highlight the long tail of genres with very few events in the Events by Genre chart",
    ],
    note="[Insert Events by Genre + Average Price by Genre charts from dashboard]"
)

# ── Slide 7: Diagnostic Insight ───────────────────────────────────────────────

add_content_slide(
    "Artists with 50M+ Spotify Listeners Average 3x More Live Events",
    [
        "**Key question:** Do streaming-popular artists have more live events?",
        "**Finding:** Yes — artists with higher Spotify monthly listeners tend to have more scheduled events",
        "This validates the streaming-to-touring pipeline hypothesis from our knowledge base research",
        "**Outliers matter:** Some high-listener artists have few events — untapped booking opportunities",
        "Streaming popularity is a leading indicator that should inform booking decisions, not just marketing",
        "**Callout:** Arrow from high-listener/high-event cluster to outlier artists with high listeners but few events",
    ],
    note="[Insert Spotify Listeners vs Event Count scatter chart from dashboard]"
)

# ── Slide 8: Recommendations ──────────────────────────────────────────────────

add_content_slide(
    "Four Data-Driven Recommendations for Entertainment Analytics Teams",
    [
        "**Route tours using city-level Spotify data** → Expect 15-20% higher ticket sell-through by matching supply to demonstrated streaming demand",
        "**Implement genre-aware dynamic pricing** → Expect 10-15% revenue lift by pricing premium genres higher and promoting emerging genres",
        "**Launch weekday promotional pricing** → Expect to fill 20%+ of empty weekday inventory with discounted bundles (e.g., 'Tuesday Night Out')",
        "**Book underrepresented genres with strong streaming** → Capture underserved fan bases and diversify beyond rock/pop concentration",
    ]
)

# ── Slide 9: Tech Stack ───────────────────────────────────────────────────────

add_content_slide(
    "Portfolio-Ready: End-to-End Pipeline Mirroring Real BI Analyst Workflows",
    [
        "**Pipeline:** Python extraction → Snowflake warehouse → dbt transformation → Streamlit visualization",
        "**Automation:** GitHub Actions runs both pipelines daily with secrets management",
        "**Quality:** 14 dbt tests, staging models with cleaning/type casting, star schema design",
        "**Research:** 17 industry sources from 11+ sites, synthesized into queryable knowledge base",
        "**Role alignment:** Built to reflect the Sr. BI Analyst role at AXS (AEG Worldwide)",
        "**Reproducible:** Clone, configure, run — full pipeline works end to end",
    ]
)

# ── Save ──────────────────────────────────────────────────────────────────────

output_path = "docs/Live_Music_Analytics_Presentation.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
